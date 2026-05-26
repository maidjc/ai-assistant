"""
AI小助手 v3.3
功能：智能对话 + 文件阅读 + 文本摘要 + 文案改写 + AI起名 + 日常小助手
运行：streamlit run app.py
依赖：pip install streamlit dashscope PyPDF2 python-docx openpyxl python-pptx
"""

import streamlit as st
from dashscope import Generation
from datetime import datetime
import io
import csv
import json


# ============ 配置区 ============
# 本地运行：在 .streamlit/secrets.toml 里写 API_KEY = "sk-xxx"
# 云端部署：在 Streamlit Cloud 的 Secrets 里加 API_KEY = "sk-xxx"
try:
    API_KEY = st.secrets["API_KEY"]
except:
    API_KEY = "sk-920a5247f801457baae1d449600a7f3d"  # 未配置时为空，会提示错误
# ================================

# 支持的文件类型
FILE_TYPES = [
    "txt", "md", "csv", "json", "xml", "yaml", "yml",
    "py", "js", "ts", "html", "css", "java", "c", "cpp", "h", "go", "rs",
    "sql", "sh", "bat", "log", "ini", "cfg", "toml",
    "pdf", "docx", "xlsx", "xls", "pptx",
]


# ===== 文件解析函数 =====

def read_text_file(uploaded_file):
    """读取纯文本类文件"""
    try:
        return uploaded_file.read().decode("utf-8")
    except UnicodeDecodeError:
        uploaded_file.seek(0)
        try:
            return uploaded_file.read().decode("gbk")
        except:
            return "[无法读取该文件内容]"


def read_csv_file(uploaded_file):
    """读取CSV文件，转成表格文字"""
    try:
        content = uploaded_file.read().decode("utf-8")
        reader = csv.reader(io.StringIO(content))
        rows = list(reader)
        if not rows:
            return "[CSV文件为空]"
        # 格式化成表格文字
        result = " | ".join(rows[0]) + "\n"
        result += " | ".join(["---"] * len(rows[0])) + "\n"
        for row in rows[1:]:
            result += " | ".join(row) + "\n"
        return result
    except Exception as e:
        return f"[CSV读取失败: {e}]"


def read_json_file(uploaded_file):
    """读取JSON文件"""
    try:
        content = uploaded_file.read().decode("utf-8")
        data = json.loads(content)
        return json.dumps(data, ensure_ascii=False, indent=2)
    except Exception as e:
        return f"[JSON读取失败: {e}]"


def read_pdf_file(uploaded_file):
    """读取PDF文件，提取文字"""
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(uploaded_file)
        text = ""
        for i, page in enumerate(reader.pages):
            page_text = page.extract_text()
            if page_text:
                text += f"--- 第{i+1}页 ---\n{page_text}\n\n"
        return text if text.strip() else "[PDF未能提取到文字内容，可能是扫描件或图片PDF]"
    except ImportError:
        return "[缺少PyPDF2库，请运行: pip install PyPDF2]"
    except Exception as e:
        return f"[PDF读取失败: {e}]"


def read_docx_file(uploaded_file):
    """读取Word文件"""
    try:
        from docx import Document
        doc = Document(uploaded_file)
        text = ""
        for para in doc.paragraphs:
            if para.text.strip():
                text += para.text + "\n"
        # 读取表格
        for i, table in enumerate(doc.tables):
            text += f"\n--- 表格{i+1} ---\n"
            for row in table.rows:
                text += " | ".join(cell.text for cell in row.cells) + "\n"
        return text if text.strip() else "[Word文档为空]"
    except ImportError:
        return "[缺少python-docx库，请运行: pip install python-docx]"
    except Exception as e:
        return f"[Word读取失败: {e}]"


def read_xlsx_file(uploaded_file):
    """读取Excel文件"""
    try:
        from openpyxl import load_workbook
        wb = load_workbook(uploaded_file, read_only=True)
        text = ""
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            text += f"--- 工作表: {sheet_name} ---\n"
            for row in ws.iter_rows(values_only=True):
                text += " | ".join(str(c) if c is not None else "" for c in row) + "\n"
            text += "\n"
        wb.close()
        return text if text.strip() else "[Excel文件为空]"
    except ImportError:
        return "[缺少openpyxl库，请运行: pip install openpyxl]"
    except Exception as e:
        return f"[Excel读取失败: {e}]"


def read_pptx_file(uploaded_file):
    """读取PPT文件"""
    try:
        from pptx import Presentation
        prs = Presentation(uploaded_file)
        text = ""
        for i, slide in enumerate(prs.slides):
            text += f"--- 第{i+1}页幻灯片 ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += shape.text + "\n"
            text += "\n"
        return text if text.strip() else "[PPT未提取到文字]"
    except ImportError:
        return "[缺少python-pptx库，请运行: pip install python-pptx]"
    except Exception as e:
        return f"[PPT读取失败: {e}]"


def parse_file(uploaded_file):
    """根据文件类型自动选择解析方式"""
    name = uploaded_file.name.lower()
    ext = name.rsplit(".", 1)[-1] if "." in name else ""

    # 文档类
    if ext == "pdf":
        return read_pdf_file(uploaded_file)
    elif ext == "docx":
        return read_docx_file(uploaded_file)
    elif ext in ("xlsx", "xls"):
        return read_xlsx_file(uploaded_file)
    elif ext == "pptx":
        return read_pptx_file(uploaded_file)

    # 数据类
    elif ext == "csv":
        return read_csv_file(uploaded_file)
    elif ext == "json":
        return read_json_file(uploaded_file)

    # 纯文本类（含代码）
    else:
        return read_text_file(uploaded_file)


# ===== AI调用 =====

def check_api_key():
    """检查API Key是否已配置"""
    if not API_KEY:
        st.error("⚠️ API Key 未配置！\n- 本地运行：在 `.streamlit/secrets.toml` 里写 `API_KEY = \"sk-xxx\"`\n- 云端部署：在 Streamlit Cloud 的 Settings → Secrets 里加 `API_KEY = \"sk-xxx\"`")
        return False
    return True


def ask_ai(messages, model="qwen-turbo", temperature=0.7):
    """调用千问大模型，返回完整回复"""
    if not API_KEY:
        return "⚠️ API Key 未配置，请先配置后再使用。"
    full = ""
    for chunk in Generation.call(
        model=model, api_key=API_KEY, messages=messages,
        stream=True, temperature=temperature, result_format="message"
    ):
        if chunk.output and chunk.output.choices:
            word = chunk.output.choices[0]['message']['content']
            if word:
                full = word
    return full


def ask_ai_stream(messages, model="qwen-turbo", temperature=0.7):
    """调用千问大模型，流式返回（聊天用）"""
    if not API_KEY:
        yield "⚠️ API Key 未配置，请先配置后再使用。"
        return
    for chunk in Generation.call(
        model=model, api_key=API_KEY, messages=messages,
        stream=True, temperature=temperature, result_format="message"
    ):
        if chunk.output and chunk.output.choices:
            word = chunk.output.choices[0]['message']['content']
            if word:
                yield word


SYSTEM_PROMPT = """你是「小智」，一个接地气的AI小助手。

【说话风格】
像朋友聊天，用大白话，不拽词不装逼，偶尔来点幽默。

【回答原则——必须严格遵守】
1. 直接回答用户的问题，不要绕弯子，不要反问，不要自说自话
2. 问什么答什么，不要发散到不相关的话题
3. 能一句说清的不写三句，简洁有力
4. 不知道就说不知道，别瞎编
5. 实用为主，少整虚的，给具体方案/步骤/建议
6. 当前日期{date}，涉及时间以这个为准
7. 用户上传了文件时，根据文件内容回答问题
8. 适当用Markdown让回答好看点"""


def main():
    st.set_page_config(
        page_title="小智 · AI助手",
        page_icon="🤖",
        layout="centered",
        initial_sidebar_state="expanded"
    )

    # ===== ngrok公网隧道（可选，用于分享）=====
    # 取消下面这行的注释就能启用公网链接
    # from pyngrok import ngrok; tunnel = ngrok.connect(8501); st.info(f"公网链接：{tunnel.public_url}")
    # ===========================================

    # ===== 全局样式 =====
    st.markdown("""
    <style>
    /* 隐藏多余元素 */
    #MainMenu, footer, header,
    [data-testid="stHeaderAction"],
    [data-testid="stHeader"],
    .stDeployButton,
    div[data-testid="stToolbar"] {visibility: hidden !important; height: 0 !important; display: none !important;}

    /* 侧边栏 */
    [data-testid="stSidebar"] {background: #fafafa !important; border-right: 1px solid #eee !important;}
    [data-testid="stSidebar"] > div:first-child {padding-top: 1rem !important;}

    /* 主背景 */
    .stApp {background: #f5f5f5;}

    /* 内容区 */
    .block-container {max-width: 720px; margin: 0 auto; padding: 0.5rem 1rem 6rem 1rem !important;}

    /* 侧边栏LOGO */
    .sidebar-logo {display: flex; align-items: center; gap: 10px; padding: 4px 0 16px 0;}
    .sidebar-logo-icon {
        width: 36px; height: 36px; border-radius: 10px;
        background: linear-gradient(135deg, #667eea, #764ba2);
        display: flex; align-items: center; justify-content: center;
        font-size: 18px; color: white;
    }
    .sidebar-logo-text {font-weight: 700; font-size: 17px; color: #1a1a2e;}

    /* 欢迎区 */
    .welcome-box {text-align: center; padding: 40px 16px 20px 16px;}
    .welcome-avatar {
        width: 60px; height: 60px; border-radius: 50%;
        background: linear-gradient(135deg, #667eea, #764ba2);
        display: inline-flex; align-items: center; justify-content: center;
        font-size: 26px; margin-bottom: 14px;
        box-shadow: 0 4px 16px rgba(102,126,234,0.2);
    }
    .welcome-title {font-size: 19px; font-weight: 700; color: #1a1a2e; margin-bottom: 4px;}
    .welcome-sub {font-size: 13px; color: #999;}

    /* 聊天气泡 */
    .stChatMessage {background: transparent !important; padding: 6px 0 !important; border: none !important;}
    [data-testid="stChatMessageAvatarUser"] > svg {display: none;}
    [data-testid="stChatMessageAvatarUser"] {
        width: 30px !important; height: 30px !important;
        background: #1a1a2e !important; border-radius: 50% !important;
    }
    [data-testid="stChatMessageAvatarUser"]::after {content: "我"; color: white; font-size: 11px; font-weight: 600;}
    [data-testid="stChatMessageAvatarAssistant"] > svg {display: none;}
    [data-testid="stChatMessageAvatarAssistant"] {
        width: 30px !important; height: 30px !important;
        background: linear-gradient(135deg, #667eea, #764ba2) !important;
        border-radius: 50% !important;
    }
    [data-testid="stChatMessageAvatarAssistant"]::after {content: "智"; color: white; font-size: 11px; font-weight: 600;}
    [data-testid="stChatMessageContent"] {font-size: 14px !important; line-height: 1.75 !important; color: #333 !important;}
    [data-testid="stChatMessageContent"] p {margin-bottom: 0.4em !important;}
    [data-testid="stChatMessageContent"] pre {
        border-radius: 10px !important; background: #1e1e2e !important;
        padding: 12px !important; font-size: 13px !important;
    }

    /* 输入框 */
    .stChatInput {border-top: none !important; padding-top: 0 !important;}
    .stChatInput > div > div > div {
        border-radius: 24px !important; border: 1.5px solid #e0e0e0 !important;
        background: #fff !important; box-shadow: 0 2px 8px rgba(0,0,0,0.04) !important;
    }
    .stChatInput > div > div > div:focus-within {
        border-color: #667eea !important; box-shadow: 0 2px 12px rgba(102,126,234,0.12) !important;
    }

    /* 功能卡片 */
    .func-card {
        background: #fff; border-radius: 16px; padding: 20px; margin-bottom: 12px;
        border: 1px solid #eee; box-shadow: 0 1px 4px rgba(0,0,0,0.03);
    }
    .func-card h3 {margin: 0 0 4px 0; font-size: 17px; color: #1a1a2e;}
    .func-card .desc {font-size: 13px; color: #999; margin-bottom: 14px;}

    /* 按钮 */
    .stButton > button {
        border-radius: 10px !important; border: 1.5px solid #e0e0e0 !important;
        background: #fff !important; color: #333 !important;
        font-size: 14px !important; font-weight: 500 !important; transition: all 0.15s !important;
    }
    .stButton > button:hover {border-color: #667eea !important; color: #667eea !important; background: #f8f7ff !important;}
    .stButton > button[kind="primary"] {background: #1a1a2e !important; color: #fff !important; border: none !important;}
    .stButton > button[kind="primary"]:hover {background: #2d2d4e !important;}

    /* 侧边栏按钮 */
    [data-testid="stSidebar"] .stButton > button {
        border: none !important; background: transparent !important;
        color: #555 !important; text-align: left !important;
        padding: 10px 12px !important; border-radius: 10px !important; font-weight: 500 !important;
    }
    [data-testid="stSidebar"] .stButton > button:hover {background: #eee !important; color: #1a1a2e !important;}

    /* 输入控件 */
    .stTextInput > div > div > input, .stTextArea > div > div > textarea, .stSelectbox > div > div > select {
        border-radius: 10px !important; border: 1.5px solid #e0e0e0 !important;
    }
    .stTextInput > div > div > input:focus, .stTextArea > div > div > textarea:focus {
        border-color: #667eea !important; box-shadow: 0 0 0 2px rgba(102,126,234,0.08) !important;
    }

    /* 侧边栏设置 */
    [data-testid="stSidebar"] label {font-size: 12px !important;}
    [data-testid="stSidebar"] .stSlider label {font-size: 12px !important; color: #999 !important;}
    .sidebar-divider {border: none; border-top: 1px solid #eee; margin: 12px 0;}
    </style>
    """, unsafe_allow_html=True)

    # ===== 初始化 =====
    if "messages" not in st.session_state:
        today = datetime.now().strftime("%Y年%m月%d日")
        st.session_state.messages = [
            {"role": "system", "content": SYSTEM_PROMPT.format(date=today)}
        ]
    if "stats" not in st.session_state:
        st.session_state.stats = {"对话次数": 0, "总字数": 0}
    if "current_func" not in st.session_state:
        st.session_state.current_func = "💬 对话"
    if "model" not in st.session_state:
        st.session_state.model = "qwen-turbo"
    if "temperature" not in st.session_state:
        st.session_state.temperature = 0.7

    # =============================================
    # 侧边栏
    # =============================================
    with st.sidebar:
        st.markdown("""
        <div class="sidebar-logo">
            <div class="sidebar-logo-icon">🤖</div>
            <div class="sidebar-logo-text">小智</div>
        </div>
        """, unsafe_allow_html=True)

        nav_items = ["💬 对话", "📝 摘要", "✍️ 改写", "🏷️ 起名", "📅 日常"]
        for label in nav_items:
            is_active = st.session_state.current_func == label
            if is_active:
                if st.button(label, key=f"nav_{label}", use_container_width=True, type="primary"):
                    pass
            else:
                if st.button(label, key=f"nav_{label}", use_container_width=True):
                    st.session_state.current_func = label
                    st.rerun()

        st.markdown('<hr class="sidebar-divider">', unsafe_allow_html=True)

        st.markdown('<div style="font-size:12px;color:#999;margin-bottom:4px;">设置</div>', unsafe_allow_html=True)
        model = st.selectbox("模型", ["qwen-turbo", "qwen-plus", "qwen-max"],
                             index=["qwen-turbo", "qwen-plus", "qwen-max"].index(st.session_state.model),
                             label_visibility="collapsed")
        temperature = st.slider("创造力", 0.0, 2.0, st.session_state.temperature, 0.1, label_visibility="collapsed")

        if model != st.session_state.model:
            st.session_state.model = model
        if temperature != st.session_state.temperature:
            st.session_state.temperature = temperature

        st.markdown('<hr class="sidebar-divider">', unsafe_allow_html=True)

        if st.button("🗑️ 清空对话", use_container_width=True):
            today = datetime.now().strftime("%Y年%m月%d日")
            st.session_state.messages = [
                {"role": "system", "content": SYSTEM_PROMPT.format(date=today)}
            ]
            st.rerun()

        chat_text = ""
        for msg in st.session_state.messages:
            if msg["role"] != "system" and isinstance(msg.get("content"), str):
                role = "我" if msg["role"] == "user" else "小智"
                chat_text += f"{role}: {msg['content']}\n\n"
        if chat_text.strip():
            st.download_button("📥 导出对话", chat_text, file_name="对话记录.txt", use_container_width=True)

        st.markdown(f"""
        <div style="position:absolute;bottom:16px;left:16px;right:16px;">
            <div style="display:flex;justify-content:space-between;font-size:12px;color:#bbb;">
                <span>💬 {st.session_state.stats['对话次数']}次</span>
                <span>{st.session_state.stats['总字数']}字</span>
            </div>
            <div style="font-size:10px;color:#ccc;text-align:center;margin-top:8px;">具身智能训练师 · 实训作品</div>
        </div>
        """, unsafe_allow_html=True)

    func = st.session_state.current_func

    # =============================================
    # 主内容区
    # =============================================

    # ===== 💬 智能对话 =====
    if func == "💬 对话":
        # 欢迎页（只有空对话才显示）
        is_empty_chat = all(m["role"] == "system" for m in st.session_state.messages)
        if is_empty_chat:
            st.markdown("""
            <div class="welcome-box">
                <div class="welcome-avatar">🤖</div>
                <div class="welcome-title">嗨，我是小智 👋</div>
                <div class="welcome-sub">能聊能写能干活，有啥说啥</div>
            </div>
            """, unsafe_allow_html=True)

        # 显示历史对话（已完成的所有回合）
        for msg in st.session_state.messages:
            if msg["role"] != "system":
                with st.chat_message(msg["role"]):
                    st.markdown(msg["content"] if isinstance(msg.get("content"), str) else "")

        # 输入框（支持多类型文件上传）
        prompt_result = st.chat_input("跟小智说点啥...", accept_file=True, file_type=FILE_TYPES)

        # 处理输入
        prompt_text = ""
        uploaded_files = None
        if prompt_result is not None:
            if hasattr(prompt_result, 'text'):
                prompt_text = prompt_result.text or ""
                uploaded_files = prompt_result.files if prompt_result.files else None
            elif isinstance(prompt_result, tuple):
                prompt_text, uploaded_files = prompt_result
                if prompt_text is None:
                    prompt_text = ""
            else:
                prompt_text = str(prompt_result) if prompt_result else ""

        if prompt_text or uploaded_files:
            user_text = prompt_text if prompt_text else "请根据上传的文件内容回答"

            # 解析上传的文件
            if uploaded_files:
                for uf in uploaded_files:
                    file_content = parse_file(uf)
                    if len(file_content) > 8000:
                        file_content = file_content[:8000] + "\n...（内容过长，已截断）"
                    user_text += f"\n\n---\n📎 文件「{uf.name}」内容：\n{file_content}"

            # 显示用户消息
            with st.chat_message("user"):
                st.markdown(prompt_text if prompt_text else "请根据上传的文件内容回答")

            # 加入消息历史
            st.session_state.messages.append({"role": "user", "content": user_text})

            # 生成AI回复
            with st.chat_message("assistant"):
                placeholder = st.empty()
                full = ""
                for word in ask_ai_stream(st.session_state.messages, st.session_state.model, st.session_state.temperature):
                    full = word
                    placeholder.markdown(full + "▌")
                placeholder.markdown(full)

            st.session_state.messages.append({"role": "assistant", "content": full})
            st.session_state.stats["对话次数"] += 1
            st.session_state.stats["总字数"] += len(full)

    # ===== 📝 文本摘要 =====
    elif func == "📝 摘要":
        st.markdown("""
        <div class="func-card">
            <h3>📝 文本摘要</h3>
            <div class="desc">贴一篇文章或上传文件，帮你抓重点</div>
        </div>
        """, unsafe_allow_html=True)

        # 支持上传文件或直接粘贴
        uploaded = st.file_uploader("上传文件", type=FILE_TYPES, key="summary_file")
        text = st.text_area("或直接粘贴文章", height=160, placeholder="把文章贴到这里...", label_visibility="collapsed" if uploaded else "visible")

        # 如果上传了文件，自动提取内容
        source_text = ""
        if uploaded:
            source_text = parse_file(uploaded)
            if len(source_text) > 8000:
                source_text = source_text[:8000] + "\n...（内容过长，已截断）"
            st.caption(f"📎 已读取: {uploaded.name}")
        if text.strip():
            source_text = text

        if st.button("✨ 生成摘要", type="primary", use_container_width=True) and source_text.strip():
            with st.spinner("提炼中..."):
                result = ask_ai([
                    {"role": "system", "content": "你是摘要助手。用3-5句话概括，突出重点，用加粗标关键词。说人话，别整学术腔。"},
                    {"role": "user", "content": source_text}
                ], st.session_state.model, 0.3)
                st.markdown(f'<div class="func-card">{result}</div>', unsafe_allow_html=True)

    # ===== ✍️ 文案改写 =====
    elif func == "✍️ 改写":
        st.markdown("""
        <div class="func-card">
            <h3>✍️ 文案改写</h3>
            <div class="desc">换个说法，让文案更好使</div>
        </div>
        """, unsafe_allow_html=True)
        style = st.selectbox("改成啥风格", ["更专业", "更口语", "更文艺", "更简洁", "更吸引人", "更搞笑"])
        text = st.text_area("输入文案", height=180, placeholder="输入你的文案...", label_visibility="collapsed")
        if st.button("✨ 改写", type="primary", use_container_width=True) and text.strip():
            with st.spinner("改写中..."):
                result = ask_ai([
                    {"role": "system", "content": f"你是文案高手。把以下内容改写得{style}，保持原意。只输出改写结果，别废话。"},
                    {"role": "user", "content": text}
                ], st.session_state.model, 0.8)
                st.markdown(f'<div class="func-card">{result}</div>', unsafe_allow_html=True)

    # ===== 🏷️ AI起名 =====
    elif func == "🏷️ 起名":
        st.markdown("""
        <div class="func-card">
            <h3>🏷️ AI起名</h3>
            <div class="desc">给品牌、产品、宠物起名字</div>
        </div>
        """, unsafe_allow_html=True)
        name_type = st.selectbox("起名类型", ["品牌/店铺", "产品/APP", "宠物", "小说角色", "群名/网名"])
        desc = st.text_input("描述一下", placeholder="比如：卖奶茶的店，面向年轻人，走国风路线")
        count = st.slider("起几个", 3, 10, 5)
        if st.button("✨ 起名", type="primary", use_container_width=True) and desc.strip():
            with st.spinner("想名字中..."):
                result = ask_ai([
                    {"role": "system", "content": f"你是起名高手。给{name_type}起名，要有创意、好记、有寓意。每个名字附一句话解释。别整太文艺的，要接地气。"},
                    {"role": "user", "content": f"给我起{count}个名字，要求：{desc}"}
                ], st.session_state.model, 0.9)
                st.markdown(f'<div class="func-card">{result}</div>', unsafe_allow_html=True)

    # ===== 📅 日常小助手 =====
    elif func == "📅 日常":
        st.markdown("""
        <div class="func-card">
            <h3>📅 日常小助手</h3>
            <div class="desc">生活里用得上的小工具</div>
        </div>
        """, unsafe_allow_html=True)

        helper_type = st.selectbox("选一个", ["📋 清单生成", "🍽️ 今天吃什么", "💪 健身计划", "📖 读书推荐", "🎬 电影推荐"])

        if helper_type == "📋 清单生成":
            list_type = st.text_input("什么清单？", placeholder="比如：去苏州旅游3天的行李清单")
            if st.button("✨ 生成清单", type="primary", use_container_width=True) and list_type.strip():
                with st.spinner("整理中..."):
                    result = ask_ai([
                        {"role": "system", "content": "你是清单达人。生成详细实用的清单，分类列出，用复选框格式。别漏关键的。"},
                        {"role": "user", "content": f"帮我生成：{list_type}"}
                    ], st.session_state.model, 0.5)
                    st.markdown(f'<div class="func-card">{result}</div>', unsafe_allow_html=True)

        elif helper_type == "🍽️ 今天吃什么":
            preference = st.text_input("口味/限制", placeholder="比如：不太辣、一人食、预算30")
            if st.button("✨ 推荐", type="primary", use_container_width=True):
                pref = preference if preference else "没有特别要求"
                with st.spinner("想菜中..."):
                    result = ask_ai([
                        {"role": "system", "content": "你是美食推荐官。推荐3-5道菜，每道说一句推荐理由和简单做法。别整那些难的。"},
                        {"role": "user", "content": f"推荐今天吃什么，要求：{pref}"}
                    ], st.session_state.model, 0.8)
                    st.markdown(f'<div class="func-card">{result}</div>', unsafe_allow_html=True)

        elif helper_type == "💪 健身计划":
            goal = st.text_input("你的目标", placeholder="比如：减脂、增肌、改善体态")
            if st.button("✨ 生成计划", type="primary", use_container_width=True) and goal.strip():
                with st.spinner("规划中..."):
                    result = ask_ai([
                        {"role": "system", "content": "你是健身教练。制定一周简易健身计划，适合新手，不用去健身房。用表格列出。动作要具体，组数次数写清楚。"},
                        {"role": "user", "content": f"目标：{goal}，给我一周计划"}
                    ], st.session_state.model, 0.5)
                    st.markdown(f'<div class="func-card">{result}</div>', unsafe_allow_html=True)

        elif helper_type == "📖 读书推荐":
            genre = st.text_input("喜欢什么类型", placeholder="比如：科幻、心理学、历史")
            if st.button("✨ 推荐", type="primary", use_container_width=True):
                g = genre if genre else "不限"
                with st.spinner("推荐中..."):
                    result = ask_ai([
                        {"role": "system", "content": "你是读书达人。推荐5本书，每本写一句话推荐理由和适合谁读。别推冷门的，要大众能找到的。"},
                        {"role": "user", "content": f"推荐{g}类的书"}
                    ], st.session_state.model, 0.7)
                    st.markdown(f'<div class="func-card">{result}</div>', unsafe_allow_html=True)

        elif helper_type == "🎬 电影推荐":
            mood = st.text_input("想看什么感觉的", placeholder="比如：轻松搞笑、烧脑悬疑、治愈暖心")
            if st.button("✨ 推荐", type="primary", use_container_width=True):
                m = mood if mood else "不限"
                with st.spinner("推荐中..."):
                    result = ask_ai([
                        {"role": "system", "content": "你是电影达人。推荐5部电影，每部写一句话推荐理由。优先推经典好片，别推烂片。"},
                        {"role": "user", "content": f"推荐{m}的电影"}
                    ], st.session_state.model, 0.7)
                    st.markdown(f'<div class="func-card">{result}</div>', unsafe_allow_html=True)


if __name__ == "__main__":
    main()
